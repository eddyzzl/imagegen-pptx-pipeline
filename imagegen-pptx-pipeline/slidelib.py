# -*- coding: utf-8 -*-
"""Shared native-shape helpers for rebuilding slide images as editable PPTX."""
import math
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

GREEN    = RGBColor(0x1E, 0x7E, 0x3E)
GREEN_BG = RGBColor(0xE9, 0xF3, 0xEC)
SLATE    = RGBColor(0x3A, 0x43, 0x50)
SLATE_MD = RGBColor(0x5A, 0x63, 0x70)
RED      = RGBColor(0xC8, 0x16, 0x1E)
RED_DK   = RGBColor(0x8F, 0x0E, 0x12)
RED_MID  = RGBColor(0xD6, 0x3A, 0x40)
PINK_BG  = RGBColor(0xFA, 0xE9, 0xE9)
PINK_LT  = RGBColor(0xFD, 0xF3, 0xF3)
PINK_LN  = RGBColor(0xF2, 0xC4, 0xC4)
BLUE     = RGBColor(0x1D, 0x5F, 0xC2)
BLUE_DK  = RGBColor(0x1D, 0x4E, 0x9E)
BLUE_BG  = RGBColor(0xE3, 0xEC, 0xFA)
BLUE_LN  = RGBColor(0xC7, 0xD9, 0xF2)
ORANGE   = RGBColor(0xE8, 0x89, 0x0C)
ORANGE_BG= RGBColor(0xFC, 0xEE, 0xDC)
DARK     = RGBColor(0x2B, 0x2B, 0x2B)
GRAY     = RGBColor(0x66, 0x66, 0x66)
GRAY_LT  = RGBColor(0x99, 0x99, 0x99)
LINE_GY  = RGBColor(0xC9, 0xC9, 0xC9)
SEP_GY   = RGBColor(0xEC, 0xEC, 0xEC)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
CARD_LN  = RGBColor(0xE9, 0xE2, 0xE2)
GRAY_MID = RGBColor(0x9E, 0x9E, 0x9E)

FONT = "Microsoft YaHei"


class SB:
    def __init__(self, src_w, src_h, bg):
        self.F = 7.5 / float(src_h)
        self.prs = Presentation()
        self.prs.slide_width = Inches(src_w * self.F)
        self.prs.slide_height = Inches(7.5)
        self.slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        self.slide.background.fill.solid()
        self.slide.background.fill.fore_color.rgb = bg
        self.shapes = self.slide.shapes

    def IN(self, px):
        return Inches(px * self.F)

    # ----- primitives -----
    def _setfont(self, run, size, color, bold, name=FONT):
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color
        run.font.name = name
        rPr = run._r.get_or_add_rPr()
        ea = rPr.find(qn('a:ea'))
        if ea is None:
            ea = rPr.makeelement(qn('a:ea'), {})
            rPr.append(ea)
        ea.set('typeface', name)

    def shape_(self, kind, x, y, w, h, fill=None, line=None, line_w=0.75,
               adj=None, rot=None, dash=None):
        sp = self.shapes.add_shape(kind, self.IN(x), self.IN(y), self.IN(w), self.IN(h))
        sp.shadow.inherit = False
        if fill is None:
            sp.fill.background()
        else:
            sp.fill.solid()
            sp.fill.fore_color.rgb = fill
        if line is None:
            sp.line.fill.background()
        else:
            sp.line.color.rgb = line
            sp.line.width = Pt(line_w)
            if dash:
                ln = sp.line._get_or_add_ln()
                ln.append(ln.makeelement(qn('a:prstDash'), {'val': dash}))
        if adj is not None:
            try:
                sp.adjustments[0] = adj
            except Exception:
                pass
        if rot is not None:
            sp.rotation = rot
        return sp

    def rect(self, x, y, w, h, fill, line=None, line_w=0.75):
        return self.shape_(MSO_SHAPE.RECTANGLE, x, y, w, h, fill, line, line_w)

    def rrect(self, x, y, w, h, fill, line=None, line_w=0.75, adj=0.12, dash=None):
        return self.shape_(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h, fill, line, line_w, adj=adj, dash=dash)

    def oval(self, x, y, w, h, fill, line=None, line_w=0.75):
        return self.shape_(MSO_SHAPE.OVAL, x, y, w, h, fill, line, line_w)

    def cline(self, x1, y1, x2, y2, color, weight=1.0, dash=None, arrow=False):
        ln = self.shapes.add_connector(1, self.IN(x1), self.IN(y1), self.IN(x2), self.IN(y2))
        ln.shadow.inherit = False
        ln.line.color.rgb = color
        ln.line.width = Pt(weight)
        l = ln.line._get_or_add_ln()
        if dash:
            l.append(l.makeelement(qn('a:prstDash'), {'val': dash}))
        if arrow:
            l.append(l.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'}))
        return ln

    def free(self, points_px, fill, line=None, line_w=0.75, close=True, dash=None):
        xs, ys = points_px[0]
        fb = self.shapes.build_freeform(self.IN(xs), self.IN(ys), scale=1.0)
        fb.add_line_segments([(self.IN(a), self.IN(b)) for a, b in points_px[1:]], close=close)
        sp = fb.convert_to_shape()
        sp.shadow.inherit = False
        if fill is None:
            sp.fill.background()
        else:
            sp.fill.solid()
            sp.fill.fore_color.rgb = fill
        if line is None:
            sp.line.fill.background()
        else:
            sp.line.color.rgb = line
            sp.line.width = Pt(line_w)
            if dash:
                ln = sp.line._get_or_add_ln()
                ln.append(ln.makeelement(qn('a:prstDash'), {'val': dash}))
        return sp

    def arc(self, cx, cy, rx, ry, a1, a2, color, w, arrow=False, dash=None, steps=24):
        """open arc polyline; angles in degrees, screen y-down (a measured CCW from +x)."""
        pts = []
        for i in range(steps + 1):
            a = math.radians(a1 + (a2 - a1) * i / steps)
            pts.append((cx + rx * math.cos(a), cy - ry * math.sin(a)))
        sp = self.free(pts, None, color, w, close=False, dash=dash)
        if arrow:
            (x1, y1), (x2, y2) = pts[-2], pts[-1]
            dx, dy = x2 - x1, y2 - y1
            L = math.hypot(dx, dy) or 1.0
            ux, uy = dx / L, dy / L
            px, py = -uy, ux
            s = w * 2.6 + 5
            self.free([(x2 + ux * s, y2 + uy * s),
                       (x2 + px * s * 0.62, y2 + py * s * 0.62),
                       (x2 - px * s * 0.62, y2 - py * s * 0.62)], color)
        return sp

    def pic(self, path, x, y, w, h, fit=True):
        if fit:
            from PIL import Image as _Im
            iw, ih = _Im.open(path).size
            k = min(w / iw, h / ih)
            nw, nh = iw * k, ih * k
            x += (w - nw) / 2.0
            y += (h - nh) / 2.0
            w, h = nw, nh
        p = self.shapes.add_picture(path, self.IN(x), self.IN(y), self.IN(w), self.IN(h))
        p.shadow.inherit = False
        return p

    def badge(self, cx, cy, r, color, num, size=8.5):
        self.oval(cx - r, cy - r, 2 * r, 2 * r, color, WHITE, 1.2)
        self.text(cx - r, cy - r, 2 * r, 2 * r, str(num), size=size, color=WHITE,
                  bold=True, align=PP_ALIGN.CENTER)

    def text(self, x, y, w, h, runs, size=12, color=DARK, bold=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE, line_spacing=None, wrap=False):
        tb = self.shapes.add_textbox(self.IN(x), self.IN(y), self.IN(w), self.IN(h))
        tf = tb.text_frame
        tf.word_wrap = wrap
        tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
        tf.vertical_anchor = anchor
        if isinstance(runs, str):
            paras = [[(runs, color, bold, size)]]
        elif runs and isinstance(runs[0], tuple):
            paras = [runs]
        else:
            paras = runs
        first = True
        for para in paras:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            p.alignment = align
            if line_spacing:
                p.line_spacing = line_spacing
            for rt in para:
                t, c, b = rt[0], rt[1], rt[2]
                s = rt[3] if len(rt) > 3 else size
                r = p.add_run()
                r.text = t
                self._setfont(r, s, c, b)
        return tb

    # ----- composite glyphs: PLACEHOLDER SCAFFOLDING ONLY -----
    # These approximate icons with native shapes. They do NOT match real source
    # artwork — using them for recognizable source icons produces a generic deck
    # (SKILL.md Iron Rule 1b). Extract real icons with iconcut3 instead. Keep these
    # only for throwaway placeholders or genuinely trivial primitives.
    def shield(self, x, y, w, h, fill, line=None, line_w=0.75):
        unit = [(0.5,0.0),(0.97,0.13),(1.0,0.2),(1.0,0.52),(0.92,0.72),(0.74,0.89),
                (0.5,1.0),(0.26,0.89),(0.08,0.72),(0.0,0.52),(0.0,0.2),(0.03,0.13)]
        return self.free([(x+u*w, y+v*h) for u, v in unit], fill, line, line_w)

    def dome(self, x, y, w, h, fill):
        pts = []
        for i in range(13):
            a = math.pi - math.pi * i / 12.0
            pts.append((x + w/2 + (w/2)*math.cos(a), y + h - h*math.sin(a)))
        return self.free(pts, fill)

    def person(self, cx, cy, scale, color):
        d = scale
        self.oval(cx-d/2, cy-d*1.05, d, d, color)
        self.dome(cx-d*1.05, cy+d*0.15, d*2.1, d*0.95, color)

    def people(self, cx, cy, color):
        self.person(cx+8, cy-1, 10, color)
        self.person(cx-4, cy+2, 13, color)

    def chevron(self, x, y, w, h, color):
        return self.shape_(MSO_SHAPE.CHEVRON, x, y, w, h, color, adj=0.6)

    def target(self, cx, cy, r, color, lw=1.5):
        self.oval(cx-r, cy-r, 2*r, 2*r, None, color, lw)
        self.oval(cx-r*0.5, cy-r*0.5, r, r, None, color, lw*0.8)
        self.oval(cx-r*0.16, cy-r*0.16, r*0.32, r*0.32, color)

    def clock(self, cx, cy, r, color, lw=1.5):
        self.oval(cx-r, cy-r, 2*r, 2*r, None, color, lw)
        self.cline(cx, cy-r*0.5, cx, cy, color, lw)
        self.cline(cx, cy, cx+r*0.45, cy+r*0.25, color, lw)

    def magnifier(self, cx, cy, r, color, lw=1.5):
        self.oval(cx-r, cy-r, r*1.5, r*1.5, None, color, lw)
        self.cline(cx+r*0.35, cy+r*0.35, cx+r, cy+r, color, lw*1.2)

    def buildings(self, cx, cy, color):
        self.rect(cx-15, cy-6, 13, 22, color)
        self.rect(cx, cy-14, 15, 30, color)
        for dy in (-9, -2, 5):
            self.rect(cx+4, cy+dy, 3, 3, WHITE)
            self.rect(cx+9, cy+dy, 3, 3, WHITE)

    def trend(self, cx, cy, color, lw=2.0, s=1.0):
        self.cline(cx-16*s, cy+12*s, cx-5*s, cy+1*s, color, lw)
        self.cline(cx-5*s, cy+1*s, cx+2*s, cy+7*s, color, lw)
        self.cline(cx+2*s, cy+7*s, cx+14*s, cy-7*s, color, lw)
        self.free([(cx+6*s,cy-9*s),(cx+17*s,cy-11*s),(cx+15*s,cy+1*s)], color)

    def coins(self, cx, cy, color, lw=1.5, s=1.0):
        self.oval(cx-15*s, cy+2*s, 30*s, 11*s, None, color, lw)
        self.oval(cx-15*s, cy-5*s, 30*s, 11*s, None, color, lw)
        self.oval(cx-15*s, cy-12*s, 30*s, 11*s, None, color, lw)
        self.oval(cx+1*s, cy+1*s, 15*s, 15*s, color)

    def bars(self, cx, cy, color, s=1.0):
        self.rect(cx-14*s, cy+2*s, 7*s, 12*s, color)
        self.rect(cx-4*s, cy-4*s, 7*s, 18*s, color)
        self.rect(cx+6*s, cy-10*s, 7*s, 24*s, color)

    def layers(self, cx, cy, color, s=1.0):
        for i, dy in enumerate((10, 0, -10)):
            self.free([(cx, cy+dy*s-7*s), (cx+16*s, cy+dy*s), (cx, cy+dy*s+7*s), (cx-16*s, cy+dy*s)], color)

    def gear(self, cx, cy, r, color):
        for i in range(8):
            a = i * 45
            rad = math.radians(a)
            tx = cx + r * math.cos(rad)
            ty = cy + r * math.sin(rad)
            self.shape_(MSO_SHAPE.RECTANGLE, tx - r*0.18, ty - r*0.14, r*0.36, r*0.28, color, rot=a+90)
        self.oval(cx-r*0.86, cy-r*0.86, r*1.72, r*1.72, color)
        self.oval(cx-r*0.34, cy-r*0.34, r*0.68, r*0.68, WHITE)

    def check_circle(self, cx, cy, r, color, txtsize=7):
        self.oval(cx-r, cy-r, 2*r, 2*r, color)
        self.text(cx-r, cy-r, 2*r, 2*r, "✓", size=txtsize, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    def save(self, path):
        self.prs.save(path)
